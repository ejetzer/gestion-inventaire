#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu Mar  3 11:22:01 2022

@author: emilejetzer
"""

import subprocess

import tkinter as tk
import configparser as cp

from pathlib import Path
from datetime import datetime as dt
from subprocess import run

import pptx
import pandas as pd

conf = cp.ConfigParser()
conf.read(Path(__file__).parent / 'nouveau_certificat.config')
sec = 'Émile'  # À changer selon qui utilise le script
CHEMIN = Path(conf.get(sec, 'Chemin')).expanduser()
FEUILLE = Path(conf.get(sec, 'Feuille')).expanduser()
UNOCONV_PY = conf.get(sec, 'unoconv')
MAJ = dt.fromisoformat(conf.get(sec, 'Dernière màj'))
UNOCONV = str(run(['which', 'unoconv'], capture_output=True).stdout, encoding='utf-8').strip()

def nouveau_certificat(nom: str, matricule: str, modèle: Path = CHEMIN):
    """
    Générer un nouveau certificat.

    Paramètres
    ----------
    nom : str
        Nom de l'étudiant ou chercheur.
    matricule : str
        Matricule de l'étudiant ou chercheur.

    Retourne
    -------
    None.

    """
    cert = pptx.Presentation(modèle)
    for forme in cert.slides[0].shapes:
        if forme.has_text_frame:
            for par in forme.text_frame.paragraphs:
                for ligne in par.runs:
                    if ligne.text == 'nom':
                        ligne.text = str(nom)
                    elif ligne.text == 'matricule':
                        ligne.text = str(matricule)
                    elif ligne.text.startswith('Date'):
                        date = dt.today()
                        ligne.text = f'Date: {date.year}-{date.month:02}'
    chemin = Path('~/Desktop').expanduser() / Path(f'cert_laser/{nom}.pptx')
    if not chemin.parent.exists():
        chemin.parent.mkdir()
    cert.save(chemin)
    #run(['unoconv', '-f', 'pdf', f'"{nom}.pptx"'])
    #temp.unlink()


def obtenir_certificats_à_faire(dernière_màj: dt = MAJ,
                                feuille: Path = FEUILLE):
    cadre = pd.read_excel(feuille,
                          usecols='C,D,E,G,AT,AU')
    cadre = cadre.rename({'Adresse de messagerie': 'courriel',
                          'Courriel (idéalement de Polytechnique)': 'courriel2',
                          'Nom': 'nom',
                          'Nom2': 'nom2',
                          'Matricule': 'matricule',
                          'Heure de fin': 'date'},
                         axis=1)\
                 .astype({'matricule': int}, errors='ignore')

    courriels_manquants = cadre['courriel'] == 'anonymous'
    cadre.loc[courriels_manquants, 'courriel'] = cadre.loc[courriels_manquants, 'courriel2']
    cadre.courriel = cadre.courriel.fillna(cadre.courriel2).fillna('@polymtl.ca')
    cadre.nom = cadre.nom.fillna(cadre.nom2).fillna('anonyme')
    cadre.date = cadre.date.dt.date
    cadre.matricule = cadre.matricule.fillna(0)
    
    return cadre.loc[cadre.date >= dernière_màj.date(), ['date', 'matricule', 'courriel', 'nom']]


class Fenetre(tk.Frame):
    """Interface pour l'entrée des informations d'un nouveau certificat."""

    def __init__(self, parent=None):
        """
        Créer une instance de Fenetre, pour entrer les informations d'un nouveau certificat.

        Paramètres
        ----------
        parent : tk.Tk, optional
            Élément d'interface parent. Par défaut, None.

        Retourne
        -------
        None.

        """
        super().__init__(parent)
        subprocess.Popen([UNOCONV_PY, UNOCONV, '--listener'])
        self.pack()
        self.creer()

    def creer(self):
        """
        Créer les éléments de l'interface.

        Retourne
        -------
        None.

        """
        # Nom de l'étudiant
        self.var_nom = tk.StringVar(self)
        self.etiquette_nom = tk.Label(self, text='Nom')
        self.champ_nom = tk.Entry(self, textvariable=self.var_nom)
        self.etiquette_nom.pack()
        self.champ_nom.pack()

        # Matricule
        self.var_matricule = tk.StringVar(self)
        self.etiquette_matricule = tk.Label(self, text='Matricule')
        self.champ_matricule = tk.Entry(self, textvariable=self.var_matricule)
        self.etiquette_matricule.pack()
        self.champ_matricule.pack()
        
        # Date de dernière màj
        self.var_dernière_màj = tk.StringVar(self, value=str(MAJ))
        self.etiquette_dernière_màj = tk.Label(self, text='Dernière màj')
        self.champ_dernière_màj = tk.Entry(self, textvariable=self.var_dernière_màj)
        self.etiquette_dernière_màj.pack()
        self.champ_dernière_màj.pack()

        # Exécuter le programme
        self.aller = tk.Button(self, text='Aller!', command=self.aller_fct)
        self.aller.pack()
        
        # Exécuter le programme automatique
        self.aller_go = tk.Button(self, text='Automatique...', command=self.auto_fct)
        self.aller_go.pack()

        # Fermeture de l'application
        self.quitter = tk.Button(self, text='Quitter.',
                                 command=self.quitter_fct)
        self.quitter.pack()
        
        # Création automatique
        self.master.after(10*60*1000, self.auto_fct2)

    def aller_fct(self):
        """
        Créer un certificat à partir des données entrées.

        Retourne
        -------
        None.

        """
        nom = self.var_nom.get()
        matricule = self.var_matricule.get()
        nouveau_certificat(nom, matricule)
        self.var_nom.set('')
        self.var_matricule.set('')
        for chemin in Path('res').glob('*.pptx'):
            subprocess.Popen([UNOCONV_PY, UNOCONV, '-f', 'pdf', str(chemin)])
    
    def auto_fct(self):
        dernière_màj = dt.fromisoformat(self.var_dernière_màj.get())
        à_faire = obtenir_certificats_à_faire(dernière_màj)
        for index, ligne in à_faire.iterrows():
            nouveau_certificat(ligne.nom, ligne.matricule)
        for chemin in Path('res').glob('*.pptx'):
            subprocess.Popen([UNOCONV_PY, UNOCONV, '-f', 'pdf', str(chemin)])
        self.var_dernière_màj.set(dt.now().isoformat())
    
    def auto_fct2(self, feuille: Path = FEUILLE):
        dernière_màj = dt.fromisoformat(self.var_dernière_màj.get())
        dernière_màj_fichier = dt.fromtimestamp(feuille.stat().st_mtime)
        if dernière_màj_fichier >= dernière_màj:
            à_faire = obtenir_certificats_à_faire(dernière_màj, feuille)
            for index, ligne in à_faire.iterrows():
                nouveau_certificat(ligne.nom, ligne.matricule)
        self.var_dernière_màj.set(dt.now().isoformat())
        for chemin in Path('res').glob('*.pptx'):
            subprocess.Popen([UNOCONV_PY, UNOCONV, '-f', 'pdf', str(chemin)])
        self.master.after(10*60*1000, self.auto_fct2)
    
    def quitter_fct(self):
        conf.set(sec, 'Dernière màj', self.var_dernière_màj.get())
        with open('nouveau_certificat.config', 'w') as fichier:
            conf.write(fichier)
        
        self.master.destroy()

def main():
    """
    Script principal: créer et faire rouler l'interface.

    Retourne
    -------
    None.

    """
    _ = tk.Tk()
    app = Fenetre(parent=_)
    app.mainloop()
